"""
Output generators: MP4 (ffmpeg), PPTX (python-pptx, opens in Keynote),
HTML (self-contained folder), and Photos Album (AppleScript).
"""
from __future__ import annotations

import html as _html
import json
import os
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import List


def generate_output(
    photos: List,
    output_dir: Path,
    base_name: str,
    output_format: str,
    prompt: str,
) -> Path:
    if output_format == "mp4":
        return _generate_mp4(photos, output_dir, base_name)
    if output_format == "pptx":
        return _generate_pptx(photos, output_dir, base_name, prompt)
    if output_format == "html":
        return _generate_html(photos, output_dir, base_name, prompt)
    if output_format == "album":
        return _generate_album(photos, output_dir, base_name, prompt)
    raise ValueError(f"Unknown format: {output_format}")


# ---------------------------------------------------------------- MP4

def _generate_mp4(photos, output_dir: Path, base_name: str) -> Path:
    if not shutil.which("ffmpeg"):
        raise RuntimeError(
            "ffmpeg not found. Install with:  brew install ffmpeg"
        )

    from PIL import Image, ImageOps

    TARGET = (1920, 1080)
    DURATION = 4.0  # seconds per image
    FPS = 30

    work = Path(tempfile.mkdtemp(prefix="slideshow_mp4_"))
    try:
        frames_dir = work / "frames"
        frames_dir.mkdir()

        valid: list[Path] = []
        for p in photos:
            src = p.path
            if not src or not os.path.exists(src):
                continue
            try:
                out = frames_dir / f"{len(valid):04d}.jpg"
                with Image.open(src) as img:
                    img = ImageOps.exif_transpose(img).convert("RGB")
                    img.thumbnail(TARGET, Image.LANCZOS)
                    canvas = Image.new("RGB", TARGET, (0, 0, 0))
                    canvas.paste(
                        img,
                        ((TARGET[0] - img.width) // 2, (TARGET[1] - img.height) // 2),
                    )
                    canvas.save(out, "JPEG", quality=90)
                valid.append(out)
            except Exception as e:
                print(f"[mp4] skipping {getattr(p, 'uuid', '?')}: {e}")

        if not valid:
            raise RuntimeError(
                "No exportable photos found for MP4 output. "
                "iCloud-only items may need to be downloaded first."
            )

        concat_file = work / "concat.txt"
        with concat_file.open("w") as f:
            for path in valid:
                f.write(f"file '{path}'\n")
                f.write(f"duration {DURATION}\n")
            # ffmpeg concat demuxer requires the last entry repeated with no duration
            f.write(f"file '{valid[-1]}'\n")

        output_path = output_dir / f"{base_name}.mp4"
        cmd = [
            "ffmpeg", "-y",
            "-f", "concat", "-safe", "0",
            "-i", str(concat_file),
            "-vsync", "vfr",
            "-vf", f"fps={FPS},format=yuv420p",
            "-c:v", "libx264",
            "-preset", "medium",
            "-crf", "20",
            str(output_path),
        ]
        result = subprocess.run(cmd, capture_output=True, text=True)
        if result.returncode != 0:
            raise RuntimeError(f"ffmpeg failed:\n{result.stderr[-2000:]}")
        return output_path
    finally:
        shutil.rmtree(work, ignore_errors=True)


# ---------------------------------------------------------------- PPTX

def _generate_pptx(photos, output_dir: Path, base_name: str, prompt: str) -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from PIL import Image as PILImage, ImageOps

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def black_bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

    # Title slide
    title_slide = prs.slides.add_slide(blank)
    black_bg(title_slide)
    tb = title_slide.shapes.add_textbox(Inches(0.75), Inches(3.0), Inches(11.83), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = prompt
    for para in tf.paragraphs:
        para.alignment = 1  # PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(36)
            run.font.color.rgb = RGBColor(255, 255, 255)

    sw, sh = prs.slide_width, prs.slide_height
    slide_ratio = sw / sh

    for p in photos:
        if not p.path or not os.path.exists(p.path):
            continue
        try:
            with PILImage.open(p.path) as img:
                img = ImageOps.exif_transpose(img)
                w, h = img.size
        except Exception:
            continue

        ratio = w / h
        if ratio > slide_ratio:
            pic_w = int(sw * 0.92)
            pic_h = int(pic_w / ratio)
        else:
            pic_h = int(sh * 0.92)
            pic_w = int(pic_h * ratio)
        left = int((sw - pic_w) / 2)
        top = int((sh - pic_h) / 2)

        slide = prs.slides.add_slide(blank)
        black_bg(slide)
        slide.shapes.add_picture(p.path, left, top, width=pic_w, height=pic_h)

    out = output_dir / f"{base_name}.pptx"
    prs.save(out)
    return out


# ---------------------------------------------------------------- HTML

def _generate_html(photos, output_dir: Path, base_name: str, prompt: str) -> Path:
    from PIL import Image, ImageOps

    assets = output_dir / f"{base_name}_assets"
    assets.mkdir(exist_ok=True)
    files: list[str] = []
    for i, p in enumerate(photos):
        if not p.path or not os.path.exists(p.path):
            continue
        try:
            dest = assets / f"{i:04d}.jpg"
            with Image.open(p.path) as img:
                img = ImageOps.exif_transpose(img).convert("RGB")
                img.thumbnail((2560, 2560), Image.LANCZOS)
                img.save(dest, "JPEG", quality=88, optimize=True)
            files.append(dest.name)
        except Exception:
            continue

    if not files:
        raise RuntimeError("No exportable photos found for HTML output.")

    html = _render_html(prompt, files, assets.name)
    out = output_dir / f"{base_name}.html"
    out.write_text(html, encoding="utf-8")
    return out


def _render_html(prompt: str, files: list[str], assets_dir: str) -> str:
    srcs = json.dumps([f"{assets_dir}/{f}" for f in files])
    title = _html.escape(prompt)
    return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>{title}</title>
<style>
  *{{margin:0;padding:0;box-sizing:border-box}}
  html,body{{height:100%;background:#000;color:#eee;font-family:-apple-system,system-ui,sans-serif;overflow:hidden}}
  #stage{{position:fixed;inset:0}}
  .slide{{position:absolute;inset:0;display:flex;align-items:center;justify-content:center;opacity:0;transition:opacity 1s ease-in-out}}
  .slide.active{{opacity:1}}
  .slide img{{max-width:96vw;max-height:96vh;object-fit:contain}}
  #title{{position:fixed;top:18px;left:20px;font-size:13px;color:rgba(255,255,255,.45);letter-spacing:.3px;z-index:5;max-width:80vw;white-space:nowrap;overflow:hidden;text-overflow:ellipsis}}
  #controls{{position:fixed;bottom:20px;left:50%;transform:translateX(-50%);display:flex;gap:10px;align-items:center;background:rgba(0,0,0,.55);padding:8px 14px;border-radius:30px;backdrop-filter:blur(10px);z-index:10;opacity:0;transition:opacity .25s}}
  body:hover #controls,body:focus-within #controls{{opacity:1}}
  button{{background:none;border:none;color:#fff;cursor:pointer;font-size:17px;padding:4px 10px;border-radius:6px}}
  button:hover{{background:rgba(255,255,255,.08)}}
  #counter{{color:rgba(255,255,255,.75);font-size:13px;padding:0 6px}}
</style>
</head>
<body>
  <div id="title">{title}</div>
  <div id="stage"></div>
  <div id="controls">
    <button id="prev" title="Previous (←)">◀</button>
    <button id="playpause" title="Play/Pause (space)">⏸</button>
    <button id="next" title="Next (→)">▶</button>
    <span id="counter">0 / 0</span>
    <button id="fs" title="Fullscreen (f)">⛶</button>
  </div>
<script>
const photos = {srcs};
const INTERVAL = 4000;
const stage = document.getElementById('stage');
const counter = document.getElementById('counter');
const pp = document.getElementById('playpause');
let idx = 0, playing = true, timer = null;

photos.forEach((src, i) => {{
  const d = document.createElement('div');
  d.className = 'slide';
  d.dataset.idx = i;
  const img = new Image();
  img.src = src;
  d.appendChild(img);
  stage.appendChild(d);
}});

function show(i) {{
  if (!photos.length) return;
  idx = ((i % photos.length) + photos.length) % photos.length;
  document.querySelectorAll('.slide').forEach(s =>
    s.classList.toggle('active', Number(s.dataset.idx) === idx));
  counter.textContent = (idx + 1) + ' / ' + photos.length;
}}
function play() {{ playing = true; pp.textContent = '⏸'; clearInterval(timer); timer = setInterval(() => show(idx + 1), INTERVAL); }}
function pause() {{ playing = false; pp.textContent = '▶'; clearInterval(timer); }}

document.getElementById('prev').onclick = () => {{ pause(); show(idx - 1); }};
document.getElementById('next').onclick = () => {{ pause(); show(idx + 1); }};
pp.onclick = () => playing ? pause() : play();
document.getElementById('fs').onclick = () =>
  document.fullscreenElement ? document.exitFullscreen() : document.documentElement.requestFullscreen();

addEventListener('keydown', e => {{
  if (e.key === 'ArrowRight') {{ pause(); show(idx + 1); }}
  else if (e.key === 'ArrowLeft') {{ pause(); show(idx - 1); }}
  else if (e.key === ' ') {{ e.preventDefault(); playing ? pause() : play(); }}
  else if (e.key === 'f' || e.key === 'F') {{ document.fullscreenElement ? document.exitFullscreen() : document.documentElement.requestFullscreen(); }}
  else if (e.key === 'Escape' && document.fullscreenElement) document.exitFullscreen();
}});

show(0); play();
</script>
</body>
</html>
"""


# ---------------------------------------------------------------- Photos album

def _generate_album(photos, output_dir: Path, base_name: str, prompt: str) -> Path:
    """
    Create an album inside Photos.app containing the selected UUIDs.
    osxphotos is read-only, so we drive Photos via AppleScript.
    """
    uuids = [p.uuid for p in photos]
    album_name = _album_title(prompt)

    script = _build_applescript(album_name, uuids)
    # Write the script to a temp file to avoid the OS argument-list length limit
    # that hits when passing a long script via osascript -e "..."
    import tempfile
    with tempfile.NamedTemporaryFile(
        mode="w", suffix=".applescript", delete=False, encoding="utf-8"
    ) as tf:
        tf.write(script)
        tf_path = tf.name
    try:
        result = subprocess.run(
            ["osascript", tf_path],
            capture_output=True, text=True,
        )
    finally:
        import os as _os
        _os.unlink(tf_path)
    if result.returncode != 0:
        raise RuntimeError(
            "Could not create album in Photos.app.\n"
            f"AppleScript error: {result.stderr.strip()}\n"
            "Make sure Terminal has automation permission for Photos "
            "(System Settings → Privacy & Security → Automation)."
        )

    report = output_dir / f"{base_name}_album.txt"
    report.write_text(
        f"Album created in Photos.app: {album_name}\n"
        f"Prompt: {prompt}\n"
        f"Photos: {len(uuids)}\n\n"
        "To play: open Photos → select this album → File → Play Slideshow.\n\n"
        "UUIDs:\n" + "\n".join(uuids),
        encoding="utf-8",
    )
    return report


def _album_title(prompt: str) -> str:
    stub = prompt.strip().replace("\n", " ")
    if len(stub) > 60:
        stub = stub[:57] + "…"
    return f"Slideshow: {stub}"


def _build_applescript(album_name: str, uuids: list[str]) -> str:
    safe_name = album_name.replace("\\", "\\\\").replace('"', '\\"')
    lines = [
        'tell application "Photos"',
        '    activate',
        f'    set newAlbum to make new album named "{safe_name}"',
    ]
    # Add each media item individually so one missing UUID doesn't abort the whole run.
    for u in uuids:
        lines.append('    try')
        lines.append(f'        set mi to media item id "{u}"')
        lines.append('        add {mi} to newAlbum')
        lines.append('    end try')
    lines.append('end tell')
    return "\n".join(lines)
